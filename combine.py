from pptx import Presentation

def extract_text_from_slide(slide):
    """Extracts all text from shapes in a slide, joined by newlines."""
    lines = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            lines.append(shape.text.strip())
    return "\n".join(lines)

def align_pptx_files(chinese_pptx_path, english_pptx_path, output_txt_path):
    prs_zh = Presentation(chinese_pptx_path)
    prs_en = Presentation(english_pptx_path)

    if len(prs_zh.slides) != len(prs_en.slides):
        raise ValueError("Slide count mismatch between Chinese and English PPTX files.")

    with open(output_txt_path, "a", encoding="utf-8") as f:
        for i, (slide_zh, slide_en) in enumerate(zip(prs_zh.slides, prs_en.slides), 1):
            zh_text = extract_text_from_slide(slide_zh).strip()
            en_text = extract_text_from_slide(slide_en).strip()

            if not zh_text or not en_text:
                continue  # skip empty slide pairs

            f.write(f"原文: {zh_text}\n")
            f.write(f"翻译: {en_text}\n\n")
            print(f"✅ Aligned Slide {i}")

    print(f"\n🎉 Done! Output saved to: {output_txt_path}")

# === Example usage ===
if __name__ == "__main__":
    output_txt = "aligned_output.txt"

    align_pptx_files("chinese\\1.pptx", "english\\1.pptx", output_txt)
    align_pptx_files("chinese\\2.pptx", "english\\2.pptx", output_txt)
    align_pptx_files("chinese\\3.pptx", "english\\3.pptx", output_txt)
    align_pptx_files("chinese\\4.pptx", "english\\4.pptx", output_txt)
    align_pptx_files("chinese\\5.pptx", "english\\5.pptx", output_txt)
    align_pptx_files("chinese\\6.pptx", "english\\6.pptx", output_txt)
    align_pptx_files("chinese\\7.pptx", "english\\7.pptx", output_txt)
    align_pptx_files("chinese\\8.pptx", "english\\8.pptx", output_txt)
    align_pptx_files("chinese\\10.pptx", "english\\10.pptx", output_txt)
    align_pptx_files("chinese\\11.pptx", "english\\11.pptx", output_txt)

