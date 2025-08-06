import os
from pptx import Presentation
from pptx.shapes.group import GroupShape
from pptx.util import Pt
from openai import OpenAI


class PowerPointTranslator:
    def __init__(self, api_key):
        self.client = OpenAI(api_key=api_key)

    def translate_batch(self, texts):
    # Assign unique IDs
        tagged_texts = [f"ID{i:03}: {t}" for i, t in enumerate(texts)]
    
        prompt = (
            "Translate the following Chinese technical phrases into fluent and PROFESSIONAL English for a manufacturing and technical context, and make sure to use correct technical terms. After translating, remove redundancy in the text and ensure that it is easy to read and grammatically correct for a presentation format. Ensure formality in tone.\n"
            "Avoid adding any bullets or numbering. Do not merge or split items.\n"
            "Translate each line while preserving the ID tag (e.g. 'ID001: ...') exactly.\n"
            "Glossary:\n"
            "密狗 - dongle\n"
            "一群人、一辈子、一件事 - Do one thing and do it well!\n"
            "股票代码 - SSC\n\n"
            "Here is the list:\n\n"
        )
        prompt += "\n".join(tagged_texts)

        try:
            print("⏳ Sending to GPT...")
            response = self.client.chat.completions.create(
            model="ft:gpt-4o-2024-08-06:personal:tr:C157tGhd",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
        )

            content = response.choices[0].message.content.strip()
            print("✅ GPT Response:\n", content)

            # Extract lines by ID tag
            translated_map = {}
            for line in content.splitlines():
                if line.startswith("ID"):
                    parts = line.split(":", 1)
                    if len(parts) == 2:
                        key, value = parts
                        translated_map[key.strip()] = value.strip()

            # Return in exact original order
            translations = [translated_map.get(f"ID{i:03}", "") for i in range(len(texts))]
            return translations

        except Exception as e:
            print(f"Translation batch failed: {e}")
            return [""] * len(texts)


    def set_font_and_resize(self, shape_or_para, new_text):
        try:
            if hasattr(shape_or_para, "runs"):
                for run in shape_or_para.runs:
                    run.text = ""
                run = shape_or_para.runs[0] if shape_or_para.runs else shape_or_para.add_run()
                run.text = new_text
                font = run.font
                font.name = "Arial"
                font.size = max(Pt(8), (font.size or Pt(16)) - Pt(4))

            elif hasattr(shape_or_para, "text_frame"):
                shape_or_para.text = new_text
                for para in shape_or_para.text_frame.paragraphs:
                    for run in para.runs:
                        font = run.font
                        font.name = "Arial"
                        font.size = max(Pt(8), (font.size or Pt(16)) - Pt(4))
        except Exception as e:
            print(f"⚠️ Font setting failed: {e}")

    def extract_shape_text(self, shape, ref_text_pairs):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    ref_text_pairs.append((para, txt))

        elif shape.shape_type == 19 and hasattr(shape, "table"):
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        ref_text_pairs.append((cell, txt))

        elif hasattr(shape, "has_chart") and shape.has_chart:
            chart = shape.chart
            if chart.has_title:
                tf = chart.chart_title.text_frame
                txt = tf.text.strip()
                if txt:
                    ref_text_pairs.append((tf, txt))

        elif isinstance(shape, GroupShape):
            for subshape in shape.shapes:
                self.extract_shape_text(subshape, ref_text_pairs)

    def translate_pptx(self, input_path, output_path, batch_size=10):
        prs = Presentation(input_path)
        ref_text_pairs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                self.extract_shape_text(shape, ref_text_pairs)

        print(f"📝 Found {len(ref_text_pairs)} items to translate.")

        translated = []
        for i in range(0, len(ref_text_pairs), batch_size):
            batch_texts = [txt for _, txt in ref_text_pairs[i:i + batch_size]]
            print(f"\n➡️ Translating batch {i}–{i + batch_size}")
            result = self.translate_batch(batch_texts)
            translated.extend(result)

        if len(translated) != len(ref_text_pairs):
            print("⚠️ Mismatch in translation count — skipping replacement to avoid offset.")
            return

        for (ref, _), new_text in zip(ref_text_pairs, translated):
            self.set_font_and_resize(ref, new_text)

        prs.save(output_path)
        print(f"\n✅ Translated PowerPoint saved to: {output_path}")


if __name__ == "__main__":
    with open("api.txt", "r") as f:
        api_key = f.read().strip()

    translator = PowerPointTranslator(api_key)
    translator.translate_pptx("v54z.pptx", "tr_v54z.pptx")
